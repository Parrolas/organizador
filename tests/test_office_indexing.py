"""Office Open XML extraction and search integration tests."""

from __future__ import annotations

import os
from pathlib import Path
from typing import NoReturn

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from organizador.db import Database
from organizador.extractors import extract_docx, extract_pptx, extract_xlsx
from organizador.indexer import MAX_INDEX_BYTES, DocumentIndexer
from organizador.models import FiledDocument, Subject


def _word_document(path: Path) -> None:
    document = Document()
    document.add_heading("Equações diferenciais", level=1)
    document.add_paragraph("Método de Euler e condições iniciais.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Passo"
    table.rows[0].cells[1].text = "Aproximação"
    document.save(path)


def _presentation(path: Path) -> None:
    presentation = Presentation()
    for title, body in (
        ("Introdução", "Referencial cartesiano"),
        ("Cinemática vetorial", "Velocidade e aceleração"),
    ):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_shape = slide.shapes.title
        assert title_shape is not None
        title_shape.text = title
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        textbox.text_frame.text = body
    presentation.save(path)


def _workbook(path: Path) -> None:
    workbook = Workbook()
    first = workbook.active
    first.title = "Medições"
    first.append(("Tempo", "Distância"))
    first.append((1, 3.5))
    second = workbook.create_sheet("Resultados")
    second.append(("Velocidade média", "=3.5/1"))
    workbook.save(path)


def _record_file(
    database: Database, subject: Subject, path: Path, kind: str = "Outros"
) -> FiledDocument:
    original = path.parent / "Downloads" / path.name
    item = database.add_inbox_item(path, original, path.name, path.stat().st_size)
    return database.record_filing(item.id, subject.id, kind, path)


def test_docx_extraction_includes_paragraphs_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "equacoes.docx"
    _word_document(path)

    pages = extract_docx(path)

    assert len(pages) == 1
    assert "Método de Euler" in pages[0]
    assert "Passo\tAproximação" in pages[0]


def test_pptx_and_xlsx_extraction_preserves_slide_and_sheet_numbers(tmp_path: Path) -> None:
    presentation_path = tmp_path / "cinematica.pptx"
    workbook_path = tmp_path / "laboratorio.xlsx"
    _presentation(presentation_path)
    _workbook(workbook_path)

    slides = extract_pptx(presentation_path)
    sheets = extract_xlsx(workbook_path)

    assert len(slides) == 2
    assert "Cinemática vetorial" in slides[1]
    assert len(sheets) == 2
    assert sheets[1].startswith("Resultados\n")
    assert "=3.5/1" in sheets[1]


def test_office_documents_are_searchable(
    database: Database, subject: Subject, tmp_path: Path
) -> None:
    word_path = tmp_path / "equacoes.docx"
    presentation_path = tmp_path / "cinematica.pptx"
    workbook_path = tmp_path / "laboratorio.xlsx"
    _word_document(word_path)
    _presentation(presentation_path)
    _workbook(workbook_path)
    documents = (
        _record_file(database, subject, word_path, "Trabalhos"),
        _record_file(database, subject, presentation_path, "Slides"),
        _record_file(database, subject, workbook_path, "Exercícios"),
    )
    indexer = DocumentIndexer(database)

    for document in documents:
        indexer.index_document(document)
    indexer.shutdown()

    assert database.search("Euler")[0].file_id == documents[0].id
    slide_result = database.search("vetorial")[0]
    assert slide_result.file_id == documents[1].id
    assert slide_result.page == 2
    sheet_result = database.search("Resultados")[0]
    assert sheet_result.file_id == documents[2].id
    assert sheet_result.page == 2


def test_corrupt_office_document_is_marked_handled(
    database: Database, subject: Subject, tmp_path: Path
) -> None:
    path = tmp_path / "corrompido.docx"
    path.write_bytes(b"not an OOXML archive")
    document = _record_file(database, subject, path)
    indexer = DocumentIndexer(database)

    indexer.index_document(document)
    indexer.shutdown()

    refreshed = database.get_file(document.id)
    assert refreshed is not None
    assert refreshed.indexed_at is not None


def test_oversized_office_document_is_not_loaded(
    database: Database,
    subject: Subject,
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    path = tmp_path / "grande.xlsx"
    path.write_bytes(b"not parsed because the on-disk size is over the limit")
    document = _record_file(database, subject, path)
    database.refresh_file_size(
        document.id,
        MAX_INDEX_BYTES + 1,
        expected_path=path,
        expected_record_token=document.record_token,
    )
    document = database.get_file(document.id)
    assert document is not None
    real_stat = Path.stat

    def fake_stat(self: Path, *, follow_symlinks: bool = True) -> os.stat_result:
        details = real_stat(self, follow_symlinks=follow_symlinks)
        if self == path:
            values = list(details)
            values[6] = MAX_INDEX_BYTES + 1
            return os.stat_result(values)
        return details

    def fail_on_parse(*args: object, **kwargs: object) -> NoReturn:
        raise AssertionError("oversized documents must not be parsed")

    monkeypatch.setattr(Path, "stat", fake_stat)
    monkeypatch.setattr("organizador.extractors.load_workbook", fail_on_parse)
    indexer = DocumentIndexer(database)

    indexer.index_document(document)
    indexer.shutdown()

    refreshed = database.get_file(document.id)
    assert refreshed is not None
    assert refreshed.indexed_at is not None
