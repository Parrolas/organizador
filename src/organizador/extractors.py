"""Text extraction for modern Office Open XML study documents."""

from __future__ import annotations

from datetime import date, datetime
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

OFFICE_SUFFIXES = frozenset({".docx", ".pptx", ".xlsx"})


def extract_docx(path: Path) -> list[str]:
    """Extract paragraphs and table cells from a Word document."""

    document = Document(str(path))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append("\t".join(cells))
    return ["\n".join(parts)]


def extract_pptx(path: Path) -> list[str]:
    """Extract text and tables while preserving PowerPoint slide numbers."""

    presentation = Presentation(str(path))
    pages: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
        pages.append("\n".join(parts))
    return pages


def extract_xlsx(path: Path) -> list[str]:
    """Extract cell values while preserving Excel worksheet numbers."""

    workbook = load_workbook(str(path), read_only=True, data_only=False, keep_links=False)
    try:
        pages: list[str] = []
        for worksheet in workbook.worksheets:
            rows = [worksheet.title]
            for row in worksheet.iter_rows(values_only=True):
                values = [_cell_text(value) for value in row]
                if any(values):
                    rows.append("\t".join(values).rstrip())
            pages.append("\n".join(rows))
        return pages
    finally:
        workbook.close()


def _cell_text(value: object) -> str:
    if value is None:
        return ""
    if isinstance(value, (date, datetime)):
        return value.isoformat()
    return str(value)
